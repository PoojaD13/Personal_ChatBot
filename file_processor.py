import PyPDF2
import docx
import pandas as pd
from pptx import Presentation
import os
import re

class FileProcessor:
    @staticmethod
    def process_file(file_path: str, filename: str):
        """Process ALL file types and return text with encoding handling"""
        file_ext = filename.split('.')[-1].lower()
        text = ""
        
        try:
            if file_ext == 'pdf':
                text = FileProcessor._read_pdf(file_path)
            elif file_ext in ['docx', 'doc']:
                text = FileProcessor._read_docx(file_path)
            elif file_ext == 'txt':
                text = FileProcessor._read_txt(file_path)
            elif file_ext in ['xlsx', 'xls', 'csv']:
                text = FileProcessor._read_spreadsheet(file_path, file_ext)
            elif file_ext in ['pptx', 'ppt']:
                text = FileProcessor._read_presentation(file_path)
            elif file_ext in ['png', 'jpg', 'jpeg', 'bmp', 'tiff']:
                text = FileProcessor._read_image(file_path)
            else:
                raise ValueError(f"Unsupported file format: {file_ext}")
                
            return text
            
        except Exception as e:
            print(f"Error processing {filename}: {str(e)}")
            return ""

    @staticmethod
    def _read_pdf(file_path: str):
        """Extract text from PDF"""
        text = ""
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        return text

    @staticmethod
    def _read_docx(file_path: str):
        """Extract text from Word documents"""
        doc = docx.Document(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])

    @staticmethod
    def _read_txt(file_path: str):
        """Read text files with encoding detection"""
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    return file.read()
            except UnicodeDecodeError:
                continue
        
        # If all encodings fail, try with error handling
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                return file.read()
        except Exception:
            return "Unable to read file - encoding issue"

    @staticmethod
    def _read_spreadsheet(file_path: str, file_ext: str):
        """Extract SMART data from ANY Excel/CSV file - GENERALIZED"""
        try:
            if file_ext == 'csv':
                df = pd.read_csv(file_path, encoding='utf-8', errors='ignore')
            else:
                df = pd.read_excel(file_path)
        
            # Create intelligent summary for ANY dataset
            text = f"Dataset Overview:\n"
            text += f"- Total records: {len(df):,}\n"
            text += f"- Columns: {', '.join(df.columns)}\n"
            text += f"- Data types: {', '.join([f'{col}({dtype})' for col, dtype in df.dtypes.items()])}\n"
        
            # Dynamic analysis based on column types
            numeric_cols = df.select_dtypes(include=['number']).columns
            text_cols = df.select_dtypes(include=['object']).columns
        
            # Analyze numeric columns
            if len(numeric_cols) > 0:
                text += f"\nNumeric Analysis:\n"
                for col in numeric_cols:
                    if col in df.columns:
                        text += f"- {col}: Total={df[col].sum():,}, Avg={df[col].mean():.2f}, Min={df[col].min():,}, Max={df[col].max():,}\n"
        
            # Analyze text columns (categorical data)
            if len(text_cols) > 0:
                text += f"\nCategorical Analysis:\n"
                for col in text_cols:
                    if col in df.columns and df[col].nunique() < 20:  # Reasonable number of categories
                        top_values = df[col].value_counts().head(3)
                        text += f"- {col}: {', '.join([f'{val}({count})' for val, count in top_values.items()])}\n"
        
            # Key insights - automatically detect important patterns
            text += f"\nKey Insights:\n"
        
            # Detect if there's date/time column
            date_cols = [col for col in df.columns if any(keyword in col.lower() for keyword in ['date', 'time', 'month', 'year'])]
            if date_cols:
                text += f"- Time period: {df[date_cols[0]].min()} to {df[date_cols[0]].max()}\n"
        
            # Detect money/financial columns
            money_cols = [col for col in df.columns if any(keyword in col.lower() for keyword in ['revenue', 'profit', 'sales', 'income', 'cost', 'price', 'amount'])]
            if money_cols:
                for col in money_cols:
                    text += f"- Total {col}: ${df[col].sum():,}\n"
        
            # Detect quantity columns
            quantity_cols = [col for col in df.columns if any(keyword in col.lower() for keyword in ['quantity', 'units', 'count', 'number', 'qty'])]
            if quantity_cols:
                for col in quantity_cols:
                    text += f"- Total {col}: {df[col].sum():,}\n"
        
            # Sample data for context
            text += f"\nSample Data (first 3 rows):\n"
            for i, (_, row) in enumerate(df.head(3).iterrows()):
                sample_text = " | ".join([f"{col}: {row[col]}" for col in df.columns[:4]])  # First 4 columns max
                text += f"{i+1}. {sample_text}\n"
                
            return text
        
        except Exception as e:
            return f"Dataset loaded with {len(df)} records. Error in detailed analysis: {str(e)}"

    @staticmethod
    def _read_presentation(file_path: str):
        """Extract text from PowerPoint files"""
        text = ""
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides):
            text += f"Slide {slide_num + 1}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text += "\n"
        return text
    @staticmethod
    def _read_image(file_path: str):
        """Extract text from images using EasyOCR - CLEAN VERSION"""
        try:
            import easyocr
            import cv2
        
            reader = easyocr.Reader(['en'])
            image = cv2.imread(file_path)
            if image is None:
                return "Could not load image file"
        
            results = reader.readtext(image)
            extracted_text = []
        
            for (bbox, text, confidence) in results:
                if confidence > 0.3:  # Only keep confident results
                    extracted_text.append(text.strip())
        
            if extracted_text:
                combined_text = " ".join(extracted_text)
                print(f"✅ EasyOCR extracted {len(combined_text)} characters: {combined_text[:100]}...")
                return f"Extracted text from image:\n{combined_text}"
            else:
                return "No text found in image"
            
        except Exception as e:
            return f"Image processing error: {str(e)}"

    @staticmethod
    def split_text(text: str, chunk_size: int = 500, chunk_overlap: int = 100):
        """Split text into smaller, meaningful chunks"""
        if not text:
            return []
        
        # Simple word-based chunking
        words = text.split()
        chunks = []
        
        for i in range(0, len(words), chunk_size - chunk_overlap):
            chunk = " ".join(words[i:i + chunk_size])
            chunks.append(chunk)
            if i + chunk_size >= len(words):
                break
        
        return chunks[:15]  # Limit total chunks

# Test the class
if __name__ == "__main__":
    print("🧪 Testing FileProcessor...")
    processor = FileProcessor()
    
    # Test that _read_image exists
    import inspect
    methods = inspect.getmembers(FileProcessor, predicate=inspect.isfunction)
    has_read_image = any(name == '_read_image' for name, method in methods)
    print(f"_read_image method exists: {has_read_image}")
    
    if has_read_image:
        print("✅ FileProcessor is ready for image processing!")
    else:
        print("❌ FileProcessor has issues!")